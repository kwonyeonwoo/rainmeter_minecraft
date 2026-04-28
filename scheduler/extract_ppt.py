from pptx import Presentation
import sys
sys.stdout.reconfigure(encoding='utf-8')

def extract_text(ppt_path):
    try:
        prs = Presentation(ppt_path)
        print(f"Total Slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"\n==================== Slide {i+1} ====================")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    print(shape.text)
    except Exception as e:
        print(f"Extraction Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        extract_text(sys.argv[1])
    else:
        print("Please provide pptx path.")
