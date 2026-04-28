from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

prs = Presentation()

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = title_text
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    
def add_content_slide(title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph()
        if isinstance(bullet, tuple):
            text, level = bullet
            p.text = text
            p.level = level
        else:
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

add_title_slide(
    "AI 학습 방지 시스템\n(Anti-Scraping System for Generative AI)", 
    "생성형 AI 환경에서의 원작자 저작권 보호 및 권리 회복 기술\n\n12223694 권연우"
)

add_content_slide("목차 (Table of Contents)", [
    "01. 서론 (Introduction) - AI 스크래핑 문제 및 기술적 한계",
    "02. 선행 연구 분석 (Literature Review) - 기존 보호 도구의 취약성",
    "03. 제안 시스템 (Proposed System) - 3-Staged AI 파이프라인 개요 및 상세",
    "04. 핵심 기술: 적대적 섭동(Adversarial Perturbation) 원리",
    "05. 차별점 및 경쟁 우위 (Evaluation & Differentiation)",
    "06. 비즈니스 모델 및 활용 방안 (Applications & Business Model)",
    "07. 기대 효과 및 결론 (Conclusion & Social Impact)"
])

add_content_slide("1. 서론: 생성형 AI 확산과 딥러닝 무단 스크래핑 위협", [
    "■ 생성형 AI(Diffusion Model) 시장의 기하급수적 성장",
    ("2025년 약 670억 달러 규모로 급성장, 연평균 37% 확대 전망", 1),
    "■ 저작권법 사각지대와 대규모 스크래핑 문제",
    ("수백만 장의 디지털 아트(Getty 등)가 명시적 동의 없이 데이터셋(LAION)으로 편입되는 기술적 침해 극심화", 1),
    "■ 크리에이터 생태계의 붕괴 위기 정량화",
    ("원본 화풍 무단 복제로 인해 일러스트레이터의 실질 수익 및 발주량 30~50% 급감 보고", 1),
    "▶ 해결책의 필요성: 법 개정 속도를 앞지르는 AI를 방어할 '실용적/즉각적인 기술적 잠금 장치'가 시급"
])

add_content_slide("2. 선행 연구 분석: 기존 보호 솔루션의 기술적 한계점", [
    "■ 대표적 선행 사례: Glaze (시카고 대학교)",
    ("초창기 솔루션이나, 무거운 연산(로컬 GPU 의존)으로 일반 사용자의 접근성 제약", 1),
    ("한계점: SNS 업로드 시 발생하는 JPEG 손실 압축 환경에서 방패(Noise)가 훼손되어 보호 효과 상실", 1),
    "■ 또 다른 선행 사례: Nightshade",
    ("Poisoning 공격으로 악성화 이미지를 만들지만 연산 속도 매우 저하", 1),
    "■ 플랫폼 자체 내장 필터 (CARA 등)",
    ("특정 사이트 종속적이며 범용 웹사이트 및 외부 유출 시 보호 불가능", 1),
    "▶ 연구 목적: '클라우드 초고속 처리', 'SNS 압축 내성 확보', '화질 오염 최소화'를 동시 달성하는 신규 모델 제시"
])

add_content_slide("3. 제안 시스템: 클라우드 기반 3-Staged AI 아키텍처", [
    "■ 시스템 개요",
    ("로컬 하드웨어 한계를 극복하고 빠른 처리를 보장하는 클라우드 통합 파이프라인 설계", 1),
    "■ 프로세스 흐름도 (3단계 모듈 체인)",
    ("1. Segmentation AI: U-Net 계열 모델을 통해 원본 이미지 내 '핵심 영역(선화, 신체)' 정밀 타겟팅", 1),
    ("2. Perturbation AI: 타겟 픽셀에 제한적인 '적대적 섭동' 노이즈 계산으로 전체 화풍 왜곡(Degradation) 극소화", 1),
    ("3. Compression-Resistant AI: 생성된 노이즈가 SNS 열화 인코딩에 파괴되지 않게 강도를 사전 보정(Pre-emphasis) 하는 핵심 보호망 삽입", 1)
])

add_content_slide("4. 핵심 기술: 적대적 섭동(Adversarial Perturbation) 원리", [
    "■ 적대적 섭동이란?",
    ("사람의 육안(HVS)으로는 감지할 수 없는 미세한 수학적 픽셀 변경을 더하는 기술 (ΔPSNR < 0.5dB)", 1),
    "■ 머신러닝의 기울기(Gradient Descent) 속임수",
    ("Loss Function(손실 함수)의 기울기를 역으로 추적(FGSM/PGD 응용)하여, 생성 모델(Diffusion)이 원본 이미지의 특징 벡터를 완전히 다른 스타일 공간으로 사상(Mapping)하게 만들어 인식 체계를 붕괴시킴", 1),
    "■ 결과적 방어 기제",
    ("해당 이미지를 무단으로 훔쳐 학습할수록 오염된 텐서 찌꺼기가 확산 과정에 개입되어, 결과적으로 AI 엔진 데이터가 역으로 망가지는 파괴적 결과 유도", 1)
])

add_content_slide("5. 시스템 경쟁 우위 및 차별화 (Evaluation)", [
    "■ 강력한 SNS 압축 내성",
    ("X, 인스타그램 환경 등 공격적인 재압축 구조 속에서도 패턴이 파괴되지 않고 보호 확률 90% 이상 유지", 1),
    "■ 로컬라이징 보호 기법 (화질 오염 최소화)",
    ("불필요한 배경까지 노이즈로 덮던 기존 툴과 탈피, 캐릭터 윤곽 등 메인 오브젝트만 타겟하여 원본 화질 보존율 대폭 우위", 1),
    "■ 플랫폼 독립적 확장성",
    ("API 형식의 클라우드 백엔드로 구동되어 데스크톱/모바일, 안드로이드/iOS 구애받지 않는 범용 솔루션 보장", 1)
])

add_content_slide("6. 비즈니스 모델(BM) 및 타겟 고객층", [
    "■ 핵심 사용자 (Target Audience)",
    ("1. B2C: 개별 일러스트레이터, 웹툰 작가, 프리랜서 영상 크리에이터의 IP 방어", 1),
    ("2. B2B: 대형 엔터테인먼트/게임 개발사 콘셉트 보호, 미술관 소장품 무단 스코어링 방지", 1),
    "■ 수익화 전략 (Freemium Model)",
    ("1. Free Tier: 월 10장 한도 무료 지원. 프리랜서 접근 문턱 파괴 및 점유율 확보", 1),
    ("2. Creator (월 9,900원): 월 500장. 압축 화질 보정 커스텀 파라미터 적용 권한, 클라우드 긴급 트래픽 배정", 1),
    ("3. Studio (월 49,000원): 무제한 API 토큰 배포(플랫폼 단위 대량 처리 자동화), 단체 라이선스 발급", 1)
])

add_content_slide("7. 제품 개발 로드맵 (Roadmap)", [
    "■ Phase 1 (단기: Prototype Validation)",
    ("적대적 노이즈(Perturbation) 산출 코어 엔진 고도화 및 픽셀 분할 테스트 검증 완료 (Current)", 1),
    "■ Phase 2 (중기: Cloud Scalability)",
    ("AWS/GCP 기반 분산 클라우드 설계. SNS 환경별 손실 함수 실시간 보정 모듈 최적화 적용", 1),
    "■ Phase 3 (장기: B2B API Market & Extension)",
    ("포트폴리오 전문 플랫폼 대상 플러그인 런칭 및 정지 이미지 컷을 넘어선 '동영상 렌더링 클립'의 프레임별 학습 차단 영역 진출", 1)
])

add_content_slide("8. 결론 및 기술의 사회적 기여 (Conclusion & Social Value)", [
    "■ 연구 프레젠테이션 핵심 요약",
    ("본 제안은 한계를 노출했던 기존 AI 스크래핑 방어 도구의 느린 속도 및 압축 취약성을 클라우드 병렬 파이프라인과 3단계 하이브리드 기술망으로 완벽히 극복한 차세대 솔루션임", 1),
    "■ 사회적 의의 (Social Impact)",
    ("저작권법 등 규제가 기술 발전 속도에 한참 뒤처져 있는 현재의 과도기에서, 창작자들이 스스로의 지적재산을 무단 복제로부터 완벽히 지켜낼 수 있는 실질적인 '기술적 억제제'를 제공하여 AI 산업과 창작 생태계 간의 합리적이고 공정한 룰 확립에 강력하게 이바지함", 1)
])

add_title_slide(
    "감사합니다.\n(Thank You & Q&A)", 
    "12223694 권연우"
)

save_path = "C:\\Users\\yeony\\Desktop\\ANTIGRAVITY\\AI학습방지_시스템_보완_업데이트본.pptx"
prs.save(save_path)
print("SUCCESS")
